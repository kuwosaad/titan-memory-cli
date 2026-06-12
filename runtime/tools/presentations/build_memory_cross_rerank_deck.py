from pathlib import Path
import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path("/Users/mohammadsaad/Desktop/Code/titan-karu")
OUT_PATH = ROOT / "docs" / "memory_cross_rerank_presentation.pptx"
TEMPLATE_PPTX = Path("/Users/mohammadsaad/Downloads/Copy of  Novah AI Pitch deck 02.pptx")
TEMPLATE_ROOT = ROOT / "tmp" / "pitch_template" / "ppt" / "media"
LOGO_WORDMARK = TEMPLATE_ROOT / "image3.png"
LOGO_MARK = TEMPLATE_ROOT / "image2.png"

BG = RGBColor(0, 1, 4)
WHITE = RGBColor(245, 245, 245)
MUTED = RGBColor(160, 167, 176)
CYAN = RGBColor(59, 222, 255)
BLUE = RGBColor(90, 121, 255)
GREEN = RGBColor(123, 245, 149)
YELLOW = RGBColor(242, 214, 90)
RED = RGBColor(255, 117, 117)

FONT = "Avenir Next"
MONO = "Menlo"


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    c_sld = slide.element.find(qn("p:cSld"))
    if c_sld is not None:
        c_sld.set("showMasterSp", "0")


def clear_slide(slide):
    sp_tree = slide.shapes._spTree  # pyright: ignore[reportAttributeAccessIssue]
    for shape in list(slide.shapes):
        sp_tree.remove(shape.element)


def add_logo(slide, small=False):
    if small:
        slide.shapes.add_picture(str(LOGO_WORDMARK), Inches(11.55), Inches(0.18), width=Inches(1.35))
    else:
        slide.shapes.add_picture(str(LOGO_WORDMARK), Inches(0.6), Inches(0.55), width=Inches(3.3))


def add_text(slide, left, top, width, height, text, size=28, color=WHITE, bold=False,
             italic=False, align=PP_ALIGN.LEFT, font=FONT, line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_block(slide, left, top, width, height, title, body, accent, title_size=22, body_size=18):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(8, 11, 16)
    shape.line.color.rgb = accent
    shape.line.width = Pt(2)
    add_text(slide, left + Inches(0.22), top + Inches(0.18), width - Inches(0.3), Inches(0.35), title,
             size=title_size, color=accent, bold=True)
    add_text(slide, left + Inches(0.22), top + Inches(0.55), width - Inches(0.35), height - Inches(0.7), body,
             size=body_size, color=WHITE)
    return shape


def add_line(slide, x1, y1, x2, y2, color=CYAN, width=2.25):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_metric(slide, left, top, width, height, label, value, accent, note):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(10, 14, 20)
    shape.line.color.rgb = accent
    shape.line.width = Pt(2)
    add_text(slide, left + Inches(0.22), top + Inches(0.18), width - Inches(0.3), Inches(0.3), label,
             size=18, color=MUTED, bold=True)
    add_text(slide, left + Inches(0.22), top + Inches(0.5), width - Inches(0.3), Inches(0.7), value,
             size=30, color=accent, bold=True)
    add_text(slide, left + Inches(0.22), top + Inches(1.1), width - Inches(0.3), height - Inches(1.25), note,
             size=16, color=WHITE)


def build_deck():
    shutil.copyfile(TEMPLATE_PPTX, OUT_PATH)
    prs = Presentation(str(OUT_PATH))
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    while len(prs.slides) < 10:
        prs.slides.add_slide(blank)
    for slide in prs.slides:
        clear_slide(slide)
    slides = [prs.slides[i] for i in range(10)]

    # Slide 1
    slide = slides[0]
    set_bg(slide)
    add_logo(slide)
    add_text(slide, Inches(0.75), Inches(2.1), Inches(10.6), Inches(1.2),
             "Cross-Memory\nReranking", size=30, color=WHITE, font=MONO)
    add_text(slide, Inches(0.8), Inches(4.25), Inches(9.5), Inches(1.0),
             "helping memories support each other before final retrieval", size=26,
             color=MUTED, italic=True)
    add_text(slide, Inches(0.8), Inches(6.55), Inches(5.4), Inches(0.3),
             "Titan-Karu research presentation", size=16, color=MUTED)
    slide.shapes.add_picture(str(LOGO_MARK), Inches(10.55), Inches(2.0), width=Inches(1.85))

    # Slide 2
    slide = slides[1]
    set_bg(slide)
    add_logo(slide, small=True)
    add_text(slide, Inches(0.55), Inches(0.45), Inches(5.8), Inches(1.2),
             "the retrieval problem", size=34, color=WHITE, italic=True)
    add_block(
        slide, Inches(0.65), Inches(1.65), Inches(5.35), Inches(4.7),
        "What current retrieval does",
        "- Embed the query\n- Score each memory in isolation\n- Return the nearest items\n\nGreat for literal matches. Weak for relational memory.",
        BLUE,
    )
    add_block(
        slide, Inches(6.35), Inches(1.65), Inches(6.1), Inches(4.7),
        "What breaks in practice",
        "- The query often matches a rough event\n- The useful answer is the learning behind it\n- One memory only becomes valuable because another memory gives it meaning\n\nThe system retrieves evidence, but misses the takeaway.",
        CYAN,
    )
    add_text(slide, Inches(0.9), Inches(6.55), Inches(11.6), Inches(0.45),
             "core failure mode: concrete memory wins even when abstract memory is the better answer", size=18,
             color=YELLOW, italic=True)

    # Slide 3
    slide = slides[2]
    set_bg(slide)
    add_logo(slide, small=True)
    add_text(slide, Inches(0.7), Inches(0.6), Inches(11.8), Inches(1.1),
             "if tokens can attend to tokens, can memories support memories?", size=31,
             color=WHITE, italic=True, align=PP_ALIGN.CENTER)
    add_block(slide, Inches(0.95), Inches(2.0), Inches(3.5), Inches(2.0),
              "rough memory", "A concrete event\nA symptom\nA local observation", BLUE)
    add_block(slide, Inches(4.92), Inches(2.0), Inches(3.5), Inches(2.0),
              "support link", "semantic overlap\nanchor terms\nturn proximity", GREEN)
    add_block(slide, Inches(8.88), Inches(2.0), Inches(3.5), Inches(2.0),
              "learning memory", "A rule\nA decision\nA compressed takeaway", CYAN)
    add_line(slide, Inches(4.45), Inches(3.0), Inches(4.9), Inches(3.0), color=GREEN, width=3)
    add_line(slide, Inches(8.42), Inches(3.0), Inches(8.86), Inches(3.0), color=GREEN, width=3)
    add_text(slide, Inches(2.1), Inches(4.75), Inches(9.2), Inches(1.2),
             "The goal is not global memory attention.\nThe goal is a safe second-stage reranker.",
             size=28, color=WHITE, align=PP_ALIGN.CENTER)

    # Slide 4
    slide = slides[3]
    set_bg(slide)
    add_logo(slide, small=True)
    add_text(slide, Inches(0.6), Inches(0.45), Inches(6.5), Inches(0.8),
             "the architecture hypothesis", size=34, color=WHITE)
    y = Inches(2.4)
    blocks = [
        ("1", "Retrieve top-k", "normal candidate generation", BLUE),
        ("2", "Measure support", "memory-to-memory links inside the pool", CYAN),
        ("3", "Rerank", "base relevance + support bonus", GREEN),
        ("4", "Return final brief", "use the reranked set", YELLOW),
    ]
    x_positions = [Inches(0.65), Inches(3.4), Inches(6.15), Inches(8.9)]
    for (idx, title, body, color), x in zip(blocks, x_positions):
        add_block(slide, x, y, Inches(2.4), Inches(2.0), idx, title + "\n" + body, color, title_size=26, body_size=16)
    add_line(slide, Inches(3.05), Inches(3.38), Inches(3.36), Inches(3.38), color=WHITE)
    add_line(slide, Inches(5.8), Inches(3.38), Inches(6.11), Inches(3.38), color=WHITE)
    add_line(slide, Inches(8.55), Inches(3.38), Inches(8.86), Inches(3.38), color=WHITE)
    add_text(slide, Inches(0.9), Inches(5.6), Inches(11.6), Inches(0.9),
             "important: this is a reranker, not a replacement for first-pass retrieval", size=24,
             color=RED, italic=True, align=PP_ALIGN.CENTER)

    # Slide 5
    slide = slides[4]
    set_bg(slide)
    add_logo(slide, small=True)
    add_text(slide, Inches(0.55), Inches(0.45), Inches(7.8), Inches(0.8),
             "why this fits Titan-Karu", size=34, color=WHITE)
    add_block(slide, Inches(0.75), Inches(1.6), Inches(5.0), Inches(3.7),
              "rough",
              "duplicate events happened because event_id collided across sessions\n\nThis is easier to match from query wording.",
              BLUE)
    add_block(slide, Inches(6.0), Inches(1.6), Inches(5.8), Inches(3.7),
              "learnings",
              "use session_id and event_id for dedupe\n\nThis is the better final answer.",
              CYAN)
    add_text(slide, Inches(1.2), Inches(5.75), Inches(10.7), Inches(0.9),
             'Query: "what rule should we use for dedupe?"\nBaseline likes the rough memory. Cross-memory reranking promotes the learning.',
             size=22, color=WHITE, align=PP_ALIGN.CENTER)

    # Slide 6
    slide = slides[5]
    set_bg(slide)
    add_logo(slide, small=True)
    add_text(slide, Inches(0.6), Inches(0.45), Inches(6.8), Inches(0.8),
             "how we tested it", size=34, color=WHITE)
    add_block(slide, Inches(0.75), Inches(1.65), Inches(3.8), Inches(3.9),
              "experiment 1",
              "Toy sanity check\n\n- mocked candidates\n- rough vs learning vs distractor\n- test whether support can fix the target failure mode",
              BLUE)
    add_block(slide, Inches(4.8), Inches(1.65), Inches(3.8), Inches(3.9),
              "experiment 2",
              "Isolated probe on the live corpus\n\n- SQLite opened read-only\n- real embeddings\n- weak-label cases mined from the store",
              GREEN)
    add_block(slide, Inches(8.85), Inches(1.65), Inches(3.8), Inches(3.9),
              "safety",
              "No writes to the live database\nNo repository initialization path\nStandalone benchmark script",
              YELLOW)
    add_text(slide, Inches(0.95), Inches(6.05), Inches(11.2), Inches(0.6),
             "probe script: tools/benchmarks/memory_cross_rerank_probe.py", size=18,
             color=MUTED, font=MONO, align=PP_ALIGN.CENTER)

    # Slide 7
    slide = slides[6]
    set_bg(slide)
    add_logo(slide, small=True)
    add_text(slide, Inches(0.6), Inches(0.45), Inches(6.0), Inches(0.8),
             "what the numbers say", size=34, color=WHITE)
    add_metric(slide, Inches(0.75), Inches(1.65), Inches(3.7), Inches(2.25),
               "MRR", "0.143 -> 0.267", CYAN, "large gain in ranking quality")
    add_metric(slide, Inches(4.8), Inches(1.65), Inches(3.7), Inches(2.25),
               "Hit@1", "0% -> 22%", GREEN, "gold moves to rank 1 far more often")
    add_metric(slide, Inches(8.85), Inches(1.65), Inches(3.7), Inches(2.25),
               "Hit@5", "32% -> 32%", RED, "no recall gain from reranking alone")
    add_block(slide, Inches(1.25), Inches(4.45), Inches(4.6), Inches(1.55),
              "candidate recall", "gold present in the initial pool: 17 / 50 cases = 34%", YELLOW,
              title_size=20, body_size=18)
    add_block(slide, Inches(6.2), Inches(4.45), Inches(5.2), Inches(1.55),
              "bottom line", "reranking helps with ordering once the right memory is already there", CYAN,
              title_size=20, body_size=18)

    # Slide 8
    slide = slides[7]
    set_bg(slide)
    add_logo(slide, small=True)
    add_text(slide, Inches(0.7), Inches(0.55), Inches(12.0), Inches(1.0),
             "interpretation", size=34, color=WHITE)
    add_text(slide, Inches(0.95), Inches(1.7), Inches(11.2), Inches(0.8),
             "The experiment clarified two separate retrieval problems.", size=28,
             color=WHITE, italic=True, align=PP_ALIGN.CENTER)
    add_block(slide, Inches(0.95), Inches(2.6), Inches(5.1), Inches(2.8),
              "Problem 1", "First-pass recall\n\nToo many gold learnings never enter the candidate pool.", RED)
    add_block(slide, Inches(6.3), Inches(2.6), Inches(5.1), Inches(2.8),
              "Problem 2", "Candidate ordering\n\nOnce the right memory is present, support-aware reranking often fixes the ranking.", GREEN)
    add_text(slide, Inches(1.1), Inches(6.0), Inches(11.1), Inches(0.7),
             "so the architecture should be: stronger recall first, cross-memory support second", size=25,
             color=CYAN, align=PP_ALIGN.CENTER)

    # Slide 9
    slide = slides[8]
    set_bg(slide)
    add_logo(slide, small=True)
    add_text(slide, Inches(0.6), Inches(0.45), Inches(6.2), Inches(0.8),
             "recommended next steps", size=34, color=WHITE)
    add_block(slide, Inches(0.8), Inches(1.65), Inches(3.65), Inches(4.2),
              "1. improve recall",
              "- retrieve a deeper pool\n- retrieve rough and learnings separately\n- add lexical or anchor expansion\n- try neighborhood expansion",
              BLUE)
    add_block(slide, Inches(4.85), Inches(1.65), Inches(3.65), Inches(4.2),
              "2. keep reranking",
              "- preserve support-aware second stage\n- measure conditional win rate\n- compare simple heuristics against attention-like scoring",
              GREEN)
    add_block(slide, Inches(8.9), Inches(1.65), Inches(3.65), Inches(4.2),
              "3. build a better benchmark",
              "- hand-label 20 to 30 queries\n- include same-session and cross-session cases\n- add positive and negative controls",
              YELLOW)

    # Slide 10
    slide = slides[9]
    set_bg(slide)
    add_logo(slide)
    add_text(slide, Inches(0.8), Inches(2.1), Inches(11.6), Inches(1.8),
             "build a stronger candidate set,\nthen let memories support each other", size=28,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.35), Inches(4.55), Inches(10.6), Inches(0.9),
             "cross-memory interaction looks real, but it belongs in stage two", size=24,
             color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(4.85), Inches(6.55), Inches(3.8), Inches(0.35),
             "thank you", size=18, color=MUTED, align=PP_ALIGN.CENTER, font=MONO)

    prs.save(str(OUT_PATH))
    print(f"wrote {OUT_PATH}")


if __name__ == "__main__":
    build_deck()
